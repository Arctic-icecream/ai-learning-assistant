from pathlib import Path

from pptx import Presentation


def extract_pptx_text(path: Path) -> str:
    presentation = Presentation(path)
    slides: list[str] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        blocks = [f"Slide {slide_number}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    blocks.append(text)
            elif shape.has_table:
                rows = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    blocks.append("\n".join(rows))

        if len(blocks) > 1:
            slides.append("\n\n".join(blocks))

    return "\n\n".join(slides)
