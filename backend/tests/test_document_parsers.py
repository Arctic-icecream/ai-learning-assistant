from pathlib import Path
from tempfile import TemporaryDirectory
import unittest

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from backend.app.document_parser import (
    UnsupportedDocumentTypeError,
    extract_document_text,
)


class DocumentParserTests(unittest.TestCase):
    def test_extracts_docx_paragraphs_and_tables_in_order(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            path = Path(temporary_directory) / "network-notes.docx"
            document = Document()
            document.add_heading("Routing Protocols", level=1)
            document.add_paragraph("AODV discovers routes on demand.")
            table = document.add_table(rows=2, cols=2)
            table.cell(0, 0).text = "Protocol"
            table.cell(0, 1).text = "Type"
            table.cell(1, 0).text = "AODV"
            table.cell(1, 1).text = "Reactive"
            document.save(path)

            text = extract_document_text(path, path.name)

            self.assertLess(text.index("Routing Protocols"), text.index("AODV discovers"))
            self.assertIn("Protocol\tType", text)
            self.assertIn("AODV\tReactive", text)

    def test_extracts_pptx_slides_text_and_tables_in_order(self) -> None:
        with TemporaryDirectory() as temporary_directory:
            path = Path(temporary_directory) / "security-models.pptx"
            presentation = Presentation()
            first_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            first_textbox = first_slide.shapes.add_textbox(
                Inches(1), Inches(1), Inches(8), Inches(1)
            )
            first_textbox.text = "Bell-LaPadula protects confidentiality."

            second_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            table_shape = second_slide.shapes.add_table(
                2, 2, Inches(1), Inches(1), Inches(8), Inches(2)
            )
            table_shape.table.cell(0, 0).text = "Rule"
            table_shape.table.cell(0, 1).text = "Meaning"
            table_shape.table.cell(1, 0).text = "No read up"
            table_shape.table.cell(1, 1).text = "Protect secrets"
            presentation.save(path)

            text = extract_document_text(path, path.name)

            self.assertIn("Slide 1", text)
            self.assertIn("Bell-LaPadula protects confidentiality.", text)
            self.assertIn("Slide 2", text)
            self.assertIn("Rule\tMeaning", text)
            self.assertLess(text.index("Slide 1"), text.index("Slide 2"))

    def test_rejects_unsupported_extensions(self) -> None:
        with self.assertRaises(UnsupportedDocumentTypeError):
            extract_document_text(Path("notes.txt"), "notes.txt")


if __name__ == "__main__":
    unittest.main()
